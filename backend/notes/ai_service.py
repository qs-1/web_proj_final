"""Gemini integration: turn messy raw input into structured study notes."""

import os
import tempfile
import json
import docx
from pptx import Presentation

from django.conf import settings
from google import genai
from google.genai import types


class AIServiceError(Exception):
    """Raised when note generation fails."""


# JSON schema we ask Gemini to produce. Matches `Note.generated_content`.
RESPONSE_SCHEMA = {
    "type": "object",
    "properties": {
        "title": {"type": "string"},
        "summary": {"type": "string"},
        "sections": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "heading": {"type": "string"},
                    "content": {"type": "string"},
                    "key_terms": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "term": {"type": "string"},
                                "definition": {"type": "string"},
                            },
                            "required": ["term", "definition"],
                        },
                    },
                },
                "required": ["heading", "content"],
            },
        },
        "key_takeaways": {"type": "array", "items": {"type": "string"}},
    },
    "required": ["title", "summary", "sections", "key_takeaways"],
}


def _build_prompt(raw_input: str, theme: str, title: str = "") -> str:
    theme_guidance = {
        "minimal": (
            "Aim for a clean, concise Notion-style result. Keep explanations tight "
            "and scannable. Prefer short paragraphs and crisp wording."
        ),
        "textbook": (
            "Aim for a polished modern-textbook result. Use slightly more formal, "
            "explanatory prose. Develop each section thoroughly with clear examples."
        ),
    }.get(theme, "")

    title_hint = f'The student suggested this title: "{title}".\n' if title else ""

    return f"""You are an expert study-notes editor for university students.

Take the student's messy, unstructured course content below and rewrite it into
clear, accurate, well-organized study notes.

Requirements:
1. Rewrite into clear, student-friendly language. Fix grammar and structure.
2. Organize the material into logical sections, each with a descriptive heading.
3. For each section, write an explanation in Markdown (you may use **bold**,
   *italics*, `code`, bullet lists, and numbered lists).
4. Extract important key terms with concise definitions for each section.
5. Write a short overall summary (2-4 sentences).
6. Provide 3-6 concise key takeaways.
7. Do NOT invent facts that are not supported by the input. Clarify and organize,
   do not fabricate.

{theme_guidance}

{title_hint}Return ONLY valid JSON matching the required schema.

--- RAW COURSE CONTENT START ---
{raw_input}
--- RAW COURSE CONTENT END ---
"""


def generate_notes(raw_input: str, theme: str = "minimal", title: str = "", file=None) -> dict:
    """Call Gemini and return a dict matching the generated_content schema."""
    if not raw_input and not file:
        raise AIServiceError("No input text or file provided.")

    api_key = settings.GEMINI_API_KEY
    if not api_key:
        raise AIServiceError(
            "GEMINI_API_KEY is not configured. Add it to backend/.env."
        )

    client = genai.Client(api_key=api_key)
    gemini_file = None
    temp_file_path = None

    if file:
        _, ext = os.path.splitext(file.name)
        ext = ext.lower()

        if ext == ".docx":
            try:
                doc = docx.Document(file)
                extracted = "\n".join([p.text for p in doc.paragraphs])
                raw_input = (raw_input or "") + "\n\n--- EXTRACTED DOCX CONTENT ---\n" + extracted
                file = None
            except Exception as e:
                raise AIServiceError(f"Failed to read DOCX file: {e}")
        elif ext == ".pptx":
            try:
                prs = Presentation(file)
                extracted = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted.append(shape.text)
                raw_input = (raw_input or "") + "\n\n--- EXTRACTED PPTX CONTENT ---\n" + "\n".join(extracted)
                file = None
            except Exception as e:
                raise AIServiceError(f"Failed to read PPTX file: {e}")

    try:
        contents = [_build_prompt(raw_input or "", theme, title)]

        if file:
            # Save the file to a temporary location for the SDK to upload
            _, ext = os.path.splitext(file.name)
            fd, temp_file_path = tempfile.mkstemp(suffix=ext)
            with os.fdopen(fd, 'wb') as f:
                for chunk in file.chunks():
                    f.write(chunk)
            
            gemini_file = client.files.upload(file=temp_file_path)
            contents.append(gemini_file)

        response = client.models.generate_content(
            model=settings.GEMINI_MODEL,
            contents=contents,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=RESPONSE_SCHEMA,
                temperature=0.4,
            ),
        )
    except Exception as exc:  # network / auth / quota errors
        raise AIServiceError(f"Gemini request failed: {exc}") from exc
    finally:
        # Clean up temporary local file
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        # Clean up remote file reference to save quota
        if gemini_file:
            try:
                client.files.delete(name=gemini_file.name)
            except Exception:
                pass

    text = (response.text or "").strip()
    if not text:
        raise AIServiceError("Gemini returned an empty response.")

    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        raise AIServiceError("Gemini returned invalid JSON.") from exc

    # Normalize / guard against missing keys so the frontend can rely on shape.
    data.setdefault("title", title or "Untitled Notes")
    data.setdefault("summary", "")
    data.setdefault("sections", [])
    data.setdefault("key_takeaways", [])
    for section in data["sections"]:
        section.setdefault("heading", "")
        section.setdefault("content", "")
        section.setdefault("key_terms", [])

    return data
